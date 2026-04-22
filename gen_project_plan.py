"""
================================================================================
  システム開発 プロジェクト計画書 (PowerPoint) ジェネレーター
  System Development Project Plan — PPTX Generator
================================================================================

  このスクリプトは 35 ページ構成のプロジェクト計画書 (.pptx) を一発生成します。
  章立て:
    1. プロジェクト定義     (目的 / スコープ / スケジュール / リスク / CP)
    2. プロジェクト体制     (体制図 / RACI / 承認プロセス)
    3. プロジェクト管理計画 (フェーズ / 会議体 / 進捗 / 変更 / 情報)
    4. 品質管理計画         (方針 / 体制 / プロセス / NFR)
    5. 移行計画             (リリース / 判定 / 保守 / 運用)
    6. AI活用計画 (Kiro)    (方針 / Specs・Hooks・Steering・MCP / ガバナンス)

================================================================================
  使い方 (3 ステップ)
================================================================================

  【前提】 Python 3.8 以降がインストールされていること。
          コマンドプロンプト / PowerShell / ターミナルで以下を実行。

  1) 依存ライブラリのインストール (初回のみ)
        pip install python-pptx

     ※ Windows で複数 Python がある場合は:
          py -m pip install python-pptx

  2) このファイルを任意の場所に保存
        例: C:\\Users\\<USER>\\Documents\\gen_project_plan.py

  3) 実行
        python gen_project_plan.py
     または:
        py gen_project_plan.py

  => 同じフォルダに `project_plan.pptx` が生成されます。
     PowerPoint でそのまま開けます (Microsoft 365 / 2019 以降で確認)。

================================================================================
  Copilot への指示例 (コピペして使える)
================================================================================

  ・「このスクリプトを実行してプロジェクト計画書を生成して」
  ・「PROJECT_NAME を『○○システム刷新計画』に変えて実行して」
  ・「OUTPUT_FILE を絶対パスに変えて、デスクトップに出力して」
  ・「会社名 (CONFIG の ORG_NAME) を『株式会社XX』に変えて」

  Copilot が「python-pptx が無い」と言ったら pip install python-pptx を指示する。

================================================================================
  カスタマイズ設定 (CONFIG) — ここだけ変更すれば OK
================================================================================
"""

# --- 必須ライブラリのインポート -----------------------------------------
# python-pptx が入っていない場合、わかりやすいエラーメッセージを出す。
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
except ImportError:
    import sys
    sys.stderr.write(
        "\n[ERROR] python-pptx がインストールされていません。\n"
        "       以下のコマンドでインストールしてください:\n\n"
        "           pip install python-pptx\n"
        "       (Windows で複数 Python がある場合: py -m pip install python-pptx)\n\n"
    )
    sys.exit(1)

import os
from pathlib import Path


# ============================================================================
#   CONFIG — プロジェクトに合わせてここを編集
# ============================================================================
# 表紙・フッタ等で表示されるプロジェクト情報。
PROJECT_NAME    = "プロジェクト計画書"
PROJECT_NAME_EN = "System Development Project Plan"
VERSION         = "1.0"
PROJECT_DATE    = "2026-04-22"              # YYYY-MM-DD
ORG_NAME        = "Project Management Office"
CONFIDENTIALITY = "Confidential / Project Plan v1.0"

# 出力ファイル名。絶対パスでも相対パスでも可。
# 相対パスの場合、このスクリプトと同じフォルダに出力されます。
OUTPUT_FILE     = "project_plan.pptx"

# 日本語フォント。Windows なら "Meiryo" or "Yu Gothic"、Mac なら "Hiragino Sans"。
# ここで指定したフォントが無ければ PowerPoint 側で自動代替されます。
JP_FONT         = "Meiryo"

# ============================================================================


# --- デザイントークン (色・サイズ) — 通常は変更不要 ----------------------
COLOR_PRIMARY = RGBColor(0x1F, 0x3A, 0x5F)   # 紺 (メインカラー)
COLOR_ACCENT  = RGBColor(0x2E, 0x86, 0xC1)   # 青 (アクセント)
COLOR_LIGHT   = RGBColor(0xEA, 0xF2, 0xF8)   # 薄青 (背景)
COLOR_TEXT    = RGBColor(0x22, 0x2B, 0x34)   # 本文
COLOR_GRAY    = RGBColor(0x6C, 0x75, 0x7D)   # 補足テキスト
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# スライドサイズ (16:9 ワイド)。
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Presentation オブジェクトを初期化。
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]  # 何も配置されていない白紙レイアウト


# ---- ユーティリティ ----------------------------------------------------
def set_font(run, size=14, bold=False, color=COLOR_TEXT, font=JP_FONT):
    run.font.name = font
    # East Asia font を明示
    rPr = run._r.get_or_add_rPr()
    eaFont = rPr.find(qn("a:ea"))
    if eaFont is None:
        from lxml import etree
        eaFont = etree.SubElement(rPr, qn("a:ea"))
    eaFont.set("typeface", font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_rect(slide, x, y, w, h, fill=COLOR_PRIMARY, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=COLOR_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    set_font(r, size=size, bold=bold, color=color)
    return tb


def add_page_header(slide, section_no, section_title, page_title):
    # 上部バー
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.7), fill=COLOR_PRIMARY)
    add_text(slide, Inches(0.4), Inches(0.12), Inches(8), Inches(0.5),
             f"{section_no}. {section_title}", size=14, bold=True,
             color=COLOR_WHITE, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, Inches(10.5), Inches(0.12), Inches(2.5), Inches(0.5),
             PROJECT_NAME, size=10, color=COLOR_WHITE,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    # ページタイトル
    add_text(slide, Inches(0.5), Inches(0.9), Inches(12), Inches(0.6),
             page_title, size=24, bold=True, color=COLOR_PRIMARY)
    # アクセント線
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(1.2), Emu(38100),
             fill=COLOR_ACCENT)


def add_footer(slide, page_no, total):
    """各ページ下部の機密表示とページ番号。"""
    add_text(slide, Inches(0.5), Inches(7.1), Inches(5), Inches(0.3),
             CONFIDENTIALITY, size=9, color=COLOR_GRAY)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{page_no} / {total}", size=9, color=COLOR_GRAY,
             align=PP_ALIGN.RIGHT)


def add_bullets(slide, x, y, w, h, items, size=13, line_spacing=1.35):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if isinstance(item, tuple):
            head, body = item
            r1 = p.add_run(); r1.text = f"■ {head}  "
            set_font(r1, size=size, bold=True, color=COLOR_PRIMARY)
            r2 = p.add_run(); r2.text = body
            set_font(r2, size=size, color=COLOR_TEXT)
        else:
            r = p.add_run(); r.text = f"・{item}"
            set_font(r, size=size, color=COLOR_TEXT)
    return tb


def add_table(slide, x, y, w, h, headers, rows,
              header_fill=COLOR_PRIMARY, header_color=COLOR_WHITE,
              zebra=COLOR_LIGHT, font_size=11):
    cols = len(headers); n_rows = len(rows) + 1
    tbl_shape = slide.shapes.add_table(n_rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    # ヘッダ
    for j, head in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.margin_top = cell.margin_bottom = Inches(0.04)
        tf = cell.text_frame; tf.word_wrap = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        r = tf.paragraphs[0].add_run(); r.text = head
        set_font(r, size=font_size, bold=True, color=header_color)
    # データ
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = zebra if i % 2 == 1 else COLOR_WHITE
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            tf = cell.text_frame; tf.word_wrap = True
            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
            r = tf.paragraphs[0].add_run(); r.text = str(val)
            set_font(r, size=font_size, color=COLOR_TEXT)
    return tbl


def add_box(slide, x, y, w, h, title, body_items,
            title_color=COLOR_PRIMARY, border=COLOR_ACCENT):
    # 枠
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    frame.fill.solid(); frame.fill.fore_color.rgb = COLOR_WHITE
    frame.line.color.rgb = border; frame.line.width = Pt(1.25)
    frame.shadow.inherit = False
    # タイトル帯
    add_rect(slide, x, y, w, Inches(0.4), fill=title_color)
    add_text(slide, x + Inches(0.15), y + Inches(0.05), w - Inches(0.3),
             Inches(0.3), title, size=12, bold=True, color=COLOR_WHITE,
             anchor=MSO_ANCHOR.MIDDLE)
    # 本文
    add_bullets(slide, x + Inches(0.2), y + Inches(0.55),
                w - Inches(0.4), h - Inches(0.6),
                body_items, size=11, line_spacing=1.25)


# ---- ページ定義 --------------------------------------------------------
# 章と各章のサブページ (section_no, section_title, page_title, builder)
PAGES = []  # (section_no, section_title, page_title, build_fn)


def builder(section_no, section_title, page_title):
    def deco(fn):
        PAGES.append((section_no, section_title, page_title, fn))
        return fn
    return deco


# ---------- 表紙 --------------------------------------------------------
def build_cover(slide):
    """表紙スライド。CONFIG の値を参照して描画する。"""
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_PRIMARY)
    add_rect(slide, 0, Inches(5.2), SLIDE_W, Inches(0.08), fill=COLOR_ACCENT)
    add_text(slide, Inches(0.8), Inches(2.3), Inches(12), Inches(1.2),
             PROJECT_NAME, size=54, bold=True, color=COLOR_WHITE)
    add_text(slide, Inches(0.8), Inches(3.4), Inches(12), Inches(0.7),
             PROJECT_NAME_EN, size=22, color=COLOR_LIGHT)
    add_text(slide, Inches(0.8), Inches(5.5), Inches(12), Inches(0.4),
             f"Version {VERSION}   /   {PROJECT_DATE}",
             size=14, color=COLOR_LIGHT)
    add_text(slide, Inches(0.8), Inches(6.0), Inches(12), Inches(0.4),
             ORG_NAME, size=14, color=COLOR_LIGHT)


# ---------- 目次 --------------------------------------------------------
def build_toc(slide):
    add_text(slide, Inches(0.6), Inches(0.5), Inches(12), Inches(0.7),
             "目次 / Table of Contents", size=28, bold=True, color=COLOR_PRIMARY)
    add_rect(slide, Inches(0.6), Inches(1.25), Inches(1.5), Emu(38100),
             fill=COLOR_ACCENT)
    toc = [
        ("1", "プロジェクト定義",
         "目的・背景・目標 / スコープ / マスタスケジュール / リスク / コンティンジェンシー"),
        ("2", "プロジェクト体制",
         "体制図 / 役割定義 / 要件決定者と承認プロセス"),
        ("3", "プロジェクト管理計画",
         "フェーズ / コミュニケーション / 進捗 / 課題・リスク / 変更 / 情報"),
        ("4", "品質管理計画",
         "方針 / 体制 / プロセス / システム品質管理詳細"),
        ("5", "移行計画",
         "リリース定義 / リリース判定 / 保守計画 / 業務運用計画"),
        ("6", "AI活用計画 (Kiro)",
         "方針 / Kiro 活用 (Spec・Hooks・Steering・MCP) / 開発プロセス統合 / ガバナンス"),
    ]
    y = Inches(1.55)
    for no, title, sub in toc:
        add_rect(slide, Inches(0.6), y, Inches(0.85), Inches(0.82),
                 fill=COLOR_PRIMARY)
        add_text(slide, Inches(0.6), y, Inches(0.85), Inches(0.82),
                 no, size=24, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, Inches(1.65), y + Inches(0.03), Inches(11), Inches(0.4),
                 title, size=17, bold=True, color=COLOR_PRIMARY)
        add_text(slide, Inches(1.65), y + Inches(0.43), Inches(11), Inches(0.4),
                 sub, size=10, color=COLOR_GRAY)
        y += Inches(0.95)


# ---------- 章扉ヘルパ --------------------------------------------------
def section_divider(no, title, subtitle):
    def _build(slide):
        add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_PRIMARY)
        add_rect(slide, Inches(0.8), Inches(2.6), Inches(1.8), Inches(0.1),
                 fill=COLOR_ACCENT)
        add_text(slide, Inches(0.8), Inches(1.9), Inches(3), Inches(0.8),
                 f"Section {no}", size=20, color=COLOR_LIGHT)
        add_text(slide, Inches(0.8), Inches(2.9), Inches(12), Inches(1.4),
                 title, size=48, bold=True, color=COLOR_WHITE)
        add_text(slide, Inches(0.8), Inches(4.3), Inches(12), Inches(0.6),
                 subtitle, size=16, color=COLOR_LIGHT)
    return _build


# ======================================================================
# 1. プロジェクト定義
# ======================================================================
@builder("1", "プロジェクト定義", "1.1 プロジェクトの目的 / 背景 / 目標")
def _(slide):
    # 3カラム
    col_w = Inches(4.1); col_h = Inches(5.0); y = Inches(1.9)
    add_box(slide, Inches(0.5), y, col_w, col_h, "背景", [
        "既存基幹システムの老朽化 (稼働14年)、保守切れ機器多数",
        "業務プロセスが属人化しており、拡張・改修コストが高騰",
        "経営方針としてデータドリブン経営への転換を加速",
        "DX推進計画 (2026-2028) の中核施策として位置付け",
    ])
    add_box(slide, Inches(4.7), y, col_w, col_h, "目的",
            [
                "老朽化システムの刷新と運用コストの大幅削減",
                "業務標準化による属人性排除と生産性向上",
                "データ基盤統合によるリアルタイム経営可視化",
                "API 連携前提のアーキテクチャで将来拡張性を確保",
            ], title_color=COLOR_ACCENT)
    add_box(slide, Inches(8.9), y, col_w, col_h, "目標 (KGI/KPI)", [
        ("KGI", "運用コスト 30% 削減 / リードタイム 40% 短縮"),
        ("KPI-1", "帳票自動化率 90% 以上"),
        ("KPI-2", "月次決算 5 営業日 → 3 営業日"),
        ("KPI-3", "本番障害 Sev1 年間 2 件以下"),
        ("期限", "2027-03 本稼働、以降 6 ヶ月ハイパーケア"),
    ])


@builder("1", "プロジェクト定義", "1.2 スコープ")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "本プロジェクトの In / Out を明確にし、認識齟齬を防止する。",
             size=12, color=COLOR_GRAY)
    # In / Out テーブル
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.0),
              headers=["区分", "領域", "対象範囲", "備考"],
              rows=[
                  ["IN", "業務", "販売 / 購買 / 在庫 / 会計 / 人事給与",
                   "BPR 実施対象"],
                  ["IN", "システム", "基幹刷新 / データ基盤 / BI / 連携基盤",
                   "クラウド (AWS) 前提"],
                  ["IN", "移行", "マスタ / トランザクション (過去3年)",
                   "クレンジング含む"],
                  ["OUT", "業務", "店舗POS / EC フロント",
                   "別プロジェクト (2027下期)"],
                  ["OUT", "システム", "メールサーバ / ファイル共有",
                   "現行資産を継続利用"],
                  ["OUT", "組織", "海外拠点 (北米・アジア)",
                   "Phase 2 にて計画"],
              ])
    add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(0.4),
             "【前提】予算は FY26 上期 確定分、為替変動は別途エスカレーション対象。",
             size=11, color=COLOR_GRAY)


@builder("1", "プロジェクト定義", "1.3 マスタスケジュール")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "フェーズゲート方式 (Waterfall + 部分 Agile)、主要マイルストーンを下記に示す。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(3.8),
              headers=["フェーズ", "期間", "主要成果物", "ゲート / 判定"],
              rows=[
                  ["要件定義", "2026-04 ～ 2026-07",
                   "要件定義書 / BPR 方針書",
                   "G1: 要件ベースライン承認"],
                  ["基本設計", "2026-07 ～ 2026-10",
                   "基本設計書 / IF 一覧 / NFR",
                   "G2: 設計レビュー完了"],
                  ["詳細設計/開発", "2026-10 ～ 2027-02",
                   "詳細設計書 / 単体テスト済ソース",
                   "G3: 単体完了"],
                  ["結合/総合テスト", "2027-01 ～ 2027-03",
                   "IT/ST 結果報告書",
                   "G4: 品質判定"],
                  ["受入/移行", "2027-02 ～ 2027-03",
                   "UAT 結果 / 移行リハ結果",
                   "G5: リリース判定"],
                  ["本稼働/保守", "2027-04 ～",
                   "ハイパーケア報告",
                   "G6: 定常運用移行"],
              ], font_size=11)
    add_text(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.4),
             "※ 各ゲートはステアリングコミッティ (月1) にて承認。NG の場合は是正まで次フェーズ着手禁止。",
             size=10, color=COLOR_GRAY)


@builder("1", "プロジェクト定義", "1.4 プロジェクトにおけるリスク")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "影響度 × 発生確率でスコアリング。High は SteerCo 報告対象。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["ID", "リスク", "影響", "確率", "スコア", "予防策"],
              rows=[
                  ["R-01", "要件肥大化 (スコープクリープ)",
                   "高", "高", "High",
                   "週次スコープレビュー / 変更管理委員会で承認"],
                  ["R-02", "キーパーソン離脱",
                   "高", "中", "High",
                   "ナレッジ二重化 / 代替要員の早期アサイン"],
                  ["R-03", "既存システム仕様の不透明",
                   "高", "高", "High",
                   "As-Is 調査フェーズを先行 / 有識者インタビュー強化"],
                  ["R-04", "性能要件 (NFR) 未達",
                   "中", "中", "Mid",
                   "PoC で早期検証 / 負荷試験の前倒し実施"],
                  ["R-05", "データ移行品質の不足",
                   "高", "中", "High",
                   "3 回のリハーサル / クレンジング基準を事前合意"],
                  ["R-06", "ベンダー調達遅延",
                   "中", "低", "Mid",
                   "代替 SI と相見積 / 契約条項にペナルティ明記"],
                  ["R-07", "ユーザ受入抵抗",
                   "中", "中", "Mid",
                   "チェンジマネジメント / 早期デモ / 教育計画"],
              ], font_size=10)


@builder("1", "プロジェクト定義", "1.5 コンティンジェンシープラン")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "主要リスクが顕在化した場合の発動条件 / 対応 / 復旧基準を定義する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["事象", "発動トリガ", "一次対応", "代替案 / 回復策", "意思決定者"],
              rows=[
                  ["要件確定遅延",
                   "要件承認 2 週間遅延",
                   "スコープ凍結 / 追加要件は Phase2 へ",
                   "MVP 切り出しで先行稼働",
                   "ステアリングコミッティ"],
                  ["性能未達",
                   "ST で SLA -20% 以上",
                   "性能改善 Task Force 編成",
                   "インフラ増強 / 一部機能のオフライン化",
                   "PM + アーキテクト"],
                  ["移行失敗",
                   "リハ 2 回連続 NG",
                   "切替判断の延期 (最大 1 ヶ月)",
                   "並行稼働期間を延長 / 段階リリース化",
                   "ステアリングコミッティ"],
                  ["本番障害 Sev1",
                   "業務停止 30 分以上",
                   "切戻し (旧システム再稼働)",
                   "BCP 手順に切替 / ホットライン開設",
                   "運用責任者 + CIO"],
                  ["要員離脱",
                   "キー2名以上の同時離脱",
                   "代替要員の即時アサイン",
                   "タスク再配分 / 一部外部委託",
                   "PM"],
              ], font_size=10)


# ======================================================================
# 2. プロジェクト体制
# ======================================================================
@builder("2", "プロジェクト体制", "2.1 体制図")
def _(slide):
    top_y = Inches(1.9)
    # ステアリング
    w_top = Inches(4.5); x_top = (SLIDE_W - w_top) / 2
    add_rect(slide, x_top, top_y, w_top, Inches(0.7), fill=COLOR_PRIMARY)
    add_text(slide, x_top, top_y, w_top, Inches(0.7),
             "ステアリングコミッティ (CIO / 役員)", size=14, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # PMO / PM
    w_mid = Inches(3.5); mid_y = Inches(3.0)
    x_pm  = (SLIDE_W - w_mid) / 2
    add_rect(slide, x_pm, mid_y, w_mid, Inches(0.7), fill=COLOR_ACCENT)
    add_text(slide, x_pm, mid_y, w_mid, Inches(0.7),
             "プロジェクトマネージャ (PM)", size=13, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # PMO / 品質
    add_rect(slide, Inches(0.8), mid_y, Inches(3.0), Inches(0.7), fill=COLOR_GRAY)
    add_text(slide, Inches(0.8), mid_y, Inches(3.0), Inches(0.7),
             "PMO (進捗・課題・情報統制)", size=12, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(slide, Inches(9.5), mid_y, Inches(3.0), Inches(0.7), fill=COLOR_GRAY)
    add_text(slide, Inches(9.5), mid_y, Inches(3.0), Inches(0.7),
             "品質保証 (QA)", size=12, bold=True,
             color=COLOR_WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # サブチーム (4列)
    sub_y = Inches(4.4); sub_w = Inches(2.9); gap = Inches(0.2)
    total_w = sub_w * 4 + gap * 3
    x0 = (SLIDE_W - total_w) / 2
    subs = [
        ("業務チーム", ["業務要件", "BPR", "UAT"]),
        ("アプリ開発", ["設計/実装", "IT/ST", "レビュー"]),
        ("基盤/インフラ", ["クラウド設計", "セキュリティ", "性能"]),
        ("移行/運用", ["データ移行", "切替計画", "運用設計"]),
    ]
    for i, (name, items) in enumerate(subs):
        x = x0 + (sub_w + gap) * i
        add_rect(slide, x, sub_y, sub_w, Inches(0.5), fill=COLOR_PRIMARY)
        add_text(slide, x, sub_y, sub_w, Inches(0.5), name, size=12,
                 bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(slide, x, sub_y + Inches(0.5), sub_w, Inches(1.6),
                 fill=COLOR_LIGHT, line=COLOR_ACCENT)
        add_bullets(slide, x + Inches(0.15), sub_y + Inches(0.6),
                    sub_w - Inches(0.3), Inches(1.5),
                    items, size=11, line_spacing=1.25)
    # 接続線 (簡易：矩形で縦線)
    line_color = COLOR_GRAY
    # Steer → PM
    add_rect(slide, Inches(6.666) - Emu(12700), top_y + Inches(0.7),
             Emu(25400), mid_y - (top_y + Inches(0.7)), fill=line_color)
    # PM → 横 (PMO / QA)
    bar_y = mid_y + Inches(0.35) - Emu(12700)
    add_rect(slide, Inches(0.8) + Inches(3.0), bar_y,
             x_pm - (Inches(0.8) + Inches(3.0)), Emu(25400), fill=line_color)
    add_rect(slide, x_pm + w_mid, bar_y,
             Inches(9.5) - (x_pm + w_mid), Emu(25400), fill=line_color)
    # PM → 各サブチーム
    add_rect(slide, Inches(6.666) - Emu(12700), mid_y + Inches(0.7),
             Emu(25400), sub_y - (mid_y + Inches(0.7)), fill=line_color)
    add_rect(slide, x0, sub_y - Emu(12700),
             total_w, Emu(25400), fill=line_color)


@builder("2", "プロジェクト体制", "2.2 役割定義 (RACI)")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "R: 実行責任  /  A: 最終承認責任  /  C: 協議  /  I: 情報共有",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.7),
              headers=["活動 / ロール", "SteerCo", "PM", "PMO", "業務", "開発", "基盤", "QA"],
              rows=[
                  ["要件確定",     "A", "R", "C", "R", "C", "C", "I"],
                  ["基本設計",     "I", "A", "C", "C", "R", "R", "C"],
                  ["開発 / 単体",  "I", "A", "I", "I", "R", "C", "C"],
                  ["結合/総合テスト", "I", "A", "C", "C", "R", "C", "R"],
                  ["UAT",          "I", "A", "C", "R", "C", "I", "C"],
                  ["リリース判定", "A", "R", "C", "C", "C", "C", "C"],
                  ["移行 / 切替",  "A", "R", "C", "C", "C", "R", "C"],
                  ["ハイパーケア", "I", "A", "C", "C", "R", "R", "C"],
              ], font_size=11)


@builder("2", "プロジェクト体制", "2.3 要件決定者と承認プロセス")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "要件は業務オーナ → PM → ステアリングの三段階で段階承認する。",
             size=12, color=COLOR_GRAY)
    # 要件決定者テーブル
    add_table(slide, Inches(0.5), Inches(2.3), Inches(6.0), Inches(4.3),
              headers=["領域", "要件決定者", "代理"],
              rows=[
                  ["販売 / 購買", "営業本部長", "営業企画部長"],
                  ["在庫 / 物流", "SCM 部長", "物流課長"],
                  ["会計 / 経理", "経理部長", "主計課長"],
                  ["人事 / 給与", "人事部長", "労務課長"],
                  ["情報 / 基盤", "情シス部長", "インフラ課長"],
                  ["全社共通", "CIO", "情シス部長"],
              ], font_size=11)
    # 承認プロセス図 (縦)
    flow_x = Inches(7.0); flow_w = Inches(5.8)
    steps = [
        ("1. 起案", "業務担当 → 業務オーナ"),
        ("2. 部門承認", "業務オーナ (要件決定者)"),
        ("3. PM レビュー", "影響 / スコープ / 工数確認"),
        ("4. 変更管理委員会", "影響 High は必須"),
        ("5. SteerCo 承認", "予算 / スケジュール影響あり"),
    ]
    y = Inches(2.3); h = Inches(0.65); gap = Inches(0.15)
    for i, (t, d) in enumerate(steps):
        add_rect(slide, flow_x, y, flow_w, h,
                 fill=COLOR_LIGHT, line=COLOR_ACCENT)
        add_text(slide, flow_x + Inches(0.2), y + Inches(0.05),
                 Inches(2), Inches(0.5), t, size=12, bold=True,
                 color=COLOR_PRIMARY, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, flow_x + Inches(2.0), y + Inches(0.05),
                 flow_w - Inches(2.1), Inches(0.5), d, size=11,
                 color=COLOR_TEXT, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            add_text(slide, flow_x + flow_w / 2 - Inches(0.2),
                     y + h, Inches(0.4), gap, "▼", size=10,
                     color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
        y += h + gap


# ======================================================================
# 3. プロジェクト管理計画
# ======================================================================
@builder("3", "プロジェクト管理計画", "3.1 フェーズ管理")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "各フェーズの完了定義 (DoD) と次フェーズ着手条件 (Gate 条件) を明確化。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["フェーズ", "目的", "主要タスク", "DoD (完了定義)", "Gate 承認者"],
              rows=[
                  ["要件定義", "業務・システム要件の確定",
                   "As-Is/To-Be / 業務整理 / NFR 合意",
                   "要件定義書 承認 / 未決課題 0",
                   "SteerCo"],
                  ["基本設計", "アーキテクチャ / IF の確定",
                   "方式 / 画面 / IF / DB 設計",
                   "基本設計書 承認 / NFR 試算 OK",
                   "PM + QA"],
                  ["詳細設計", "実装仕様の確定",
                   "モジュール / バッチ / IF 詳細",
                   "詳細設計レビュー 100% 完了",
                   "PM + アーキ"],
                  ["開発", "実装と単体品質の確保",
                   "コーディング / 単体テスト / レビュー",
                   "UT カバレッジ ≥ 80% / 欠陥残 0",
                   "PM"],
                  ["テスト", "結合~総合での品質確保",
                   "IT / ST / 性能 / セキュリティ",
                   "Sev1/2 欠陥 0 / SLA 達成",
                   "QA + PM"],
                  ["移行 / 本番", "安全なサービスイン",
                   "UAT / リハ / 切替",
                   "リリース判定通過 / BCP 確認",
                   "SteerCo"],
              ], font_size=10)


@builder("3", "プロジェクト管理計画", "3.2 コミュニケーション管理")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "会議体と情報経路を固定化し、意思決定のリードタイムを短縮する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["会議体", "頻度", "参加者", "目的", "主なアウトプット"],
              rows=[
                  ["ステアリングコミッティ", "月1",
                   "役員 / CIO / PM",
                   "重要意思決定 / 予算 / リスク",
                   "議事録 / 承認事項"],
                  ["PM 定例", "週1",
                   "PM / PMO / 各リーダ",
                   "進捗 / 課題 / 変更",
                   "進捗表 / 課題一覧"],
                  ["チーム定例", "週2",
                   "チームメンバ",
                   "タスク進捗 / 障害",
                   "Daily Log"],
                  ["要件レビュー", "都度",
                   "業務オーナ / 業務 / 開発",
                   "要件仕様の確定",
                   "レビュー結果 / 承認"],
                  ["障害対応会", "発生時",
                   "PM / QA / 開発 / 運用",
                   "復旧 / 再発防止",
                   "障害報告書"],
                  ["ユーザ説明会", "フェーズ毎",
                   "現場ユーザ",
                   "理解促進 / 受入準備",
                   "説明資料 / Q&A"],
              ], font_size=10)


@builder("3", "プロジェクト管理計画", "3.3 進捗管理")
def _(slide):
    col_w = Inches(6.1); y = Inches(1.9); h = Inches(5.0)
    add_box(slide, Inches(0.5), y, col_w, h, "測定と可視化", [
        ("WBS 粒度", "最小タスク ≤ 5 人日 / ≤ 10 営業日"),
        ("指標", "EVM (PV/EV/AC), SPI, CPI を週次で算出"),
        ("見える化", "Jira + 週次バーンダウン / ガントチャート"),
        ("早期検知", "SPI < 0.9 で黄信号、< 0.8 で赤信号"),
        ("成果物", "週次進捗報告 / 月次 SteerCo 報告"),
    ])
    add_box(slide, Inches(6.7), y, col_w, h, "遅延時のアクション",
            [
                ("黄信号", "チームリーダーが原因分析・PM 報告"),
                ("赤信号", "リカバリプラン策定 (1週間以内)"),
                ("手段", "タスク再配分 / 要員増員 / スコープ調整"),
                ("エスカレ", "2 週間で回復不能 → SteerCo 付議"),
                ("記録", "RCA (根本原因分析) を必ず文書化"),
            ], title_color=COLOR_ACCENT)


@builder("3", "プロジェクト管理計画", "3.4 課題 / リスク管理")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "課題 (Issue) と リスク (Risk) を分離管理。優先度でリソースを配分。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.4),
              headers=["区分", "定義", "管理ツール", "更新頻度", "レビュー"],
              rows=[
                  ["Issue (課題)", "既に発生している事象",
                   "Jira (Issue project)", "日次", "PM 定例"],
                  ["Risk (リスク)", "将来発生しうる事象",
                   "Risk Register (Confluence)", "週次", "PM 定例 + SteerCo"],
              ], font_size=11)
    # 優先度マトリクス
    y2 = Inches(4.9)
    add_text(slide, Inches(0.5), y2, Inches(12), Inches(0.4),
             "優先度マトリクス (影響 × 緊急度)", size=14, bold=True,
             color=COLOR_PRIMARY)
    add_table(slide, Inches(0.5), y2 + Inches(0.5), Inches(12.3), Inches(1.6),
              headers=["", "緊急度 低", "緊急度 中", "緊急度 高"],
              rows=[
                  ["影響 高", "P2 / 計画対応", "P1 / 即対応", "P0 / 緊急対応"],
                  ["影響 中", "P3 / 監視",    "P2 / 計画対応", "P1 / 即対応"],
                  ["影響 低", "P4 / 受容",    "P3 / 監視",    "P2 / 計画対応"],
              ], font_size=11)


@builder("3", "プロジェクト管理計画", "3.5 変更管理")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "スコープ / スケジュール / 予算 / 品質に影響する変更は CCB で審議。",
             size=12, color=COLOR_GRAY)
    # プロセス
    steps = ["変更要求 (CR)", "影響分析", "CCB 審議", "承認/却下", "計画反映", "実行/検証"]
    x = Inches(0.5); y = Inches(2.4); w = Inches(2.0); h = Inches(0.8)
    gap_x = Inches(0.12)
    for i, s in enumerate(steps):
        xi = x + (w + gap_x) * i
        add_rect(slide, xi, y, w, h,
                 fill=COLOR_ACCENT if i % 2 == 0 else COLOR_PRIMARY)
        add_text(slide, xi, y, w, h, s, size=11, bold=True,
                 color=COLOR_WHITE, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
    # 判定基準
    add_text(slide, Inches(0.5), Inches(3.6), Inches(12), Inches(0.4),
             "判定基準 (目安)", size=14, bold=True, color=COLOR_PRIMARY)
    add_table(slide, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.4),
              headers=["区分", "影響", "承認者", "SLA"],
              rows=[
                  ["軽微 (Minor)", "工数 ≤ 5 人日 / スケジュール影響なし",
                   "PM", "3 営業日以内"],
                  ["中 (Major)", "工数 5〜20 人日 / 一部機能影響",
                   "PM + 業務オーナ", "5 営業日以内"],
                  ["重大 (Critical)", "工数 > 20 人日 / スコープ・予算影響",
                   "CCB + SteerCo", "10 営業日以内"],
              ], font_size=11)


@builder("3", "プロジェクト管理計画", "3.6 情報管理")
def _(slide):
    col_w = Inches(6.1); y = Inches(1.9); h = Inches(5.0)
    add_box(slide, Inches(0.5), y, col_w, h, "情報資産の分類と格納先", [
        ("ドキュメント", "Confluence (WBS / 議事 / 設計書)"),
        ("ソース",       "GitHub Enterprise (Private)"),
        ("チケット",     "Jira (要件 / 課題 / 変更)"),
        ("成果物バイナリ", "Artifact Repo (Nexus)"),
        ("機密資料",     "暗号化共有ストレージ / IRM 必須"),
        ("チャット",     "Slack (ログ保管 180 日)"),
    ])
    add_box(slide, Inches(6.7), y, col_w, h, "アクセス制御と監査",
            [
                ("権限モデル", "ロールベース (RBAC) / 最小権限"),
                ("機密区分",   "Public / Internal / Confidential / Secret"),
                ("持出制御",   "Secret 相当は端末持出禁止 (DLP 監視)"),
                ("ログ",       "アクセス / 変更ログを1年保管"),
                ("監査",       "四半期ごとに PMO がサンプル監査"),
                ("廃棄",       "プロジェクト終了後1年で削除 (法令除く)"),
            ], title_color=COLOR_ACCENT)


# ======================================================================
# 4. 品質管理計画
# ======================================================================
@builder("4", "品質管理計画", "4.1 品質管理方針")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "「作り込みで品質を上げる」を基本方針とし、後工程での手戻りを抑制する。",
             size=12, color=COLOR_GRAY)
    items = [
        ("方針1", "品質は工程ごとに作り込む (Shift-Left)"),
        ("方針2", "欠陥は発見工程ではなく混入工程で是正する"),
        ("方針3", "定量 (メトリクス) と定性 (レビュー) の両輪で評価する"),
        ("方針4", "自動化可能な検証は CI に組み込み、人は設計レビューに集中"),
        ("方針5", "顧客品質 (利用者価値) と内部品質 (保守性) を共に重視"),
    ]
    add_bullets(slide, Inches(0.8), Inches(2.4), Inches(11.8), Inches(3),
                items, size=14, line_spacing=1.5)
    # ISO9126 / 25010 観点
    add_text(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
             "評価観点 (ISO/IEC 25010 ベース)", size=14, bold=True,
             color=COLOR_PRIMARY)
    cats = ["機能適合性", "性能効率", "互換性", "使用性", "信頼性", "セキュリティ", "保守性", "移植性"]
    x = Inches(0.5); y = Inches(5.7); w = Inches(1.55); h = Inches(0.7)
    for i, c in enumerate(cats):
        add_rect(slide, x + (w + Inches(0.05)) * i, y, w, h,
                 fill=COLOR_LIGHT, line=COLOR_ACCENT)
        add_text(slide, x + (w + Inches(0.05)) * i, y, w, h, c,
                 size=11, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@builder("4", "品質管理計画", "4.2 品質管理体制")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "QA は開発から独立したレポートラインで品質を担保する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["ロール", "所属", "主責任", "レポート先"],
              rows=[
                  ["品質統括", "QA 部",
                   "品質方針策定 / 品質判定",
                   "PM / SteerCo"],
                  ["品質管理リーダ", "プロジェクト",
                   "メトリクス収集 / 分析 / 是正指示",
                   "品質統括"],
                  ["設計レビュア", "アーキ/有識者",
                   "設計レビュー / NFR 確認",
                   "品質管理リーダ"],
                  ["テストリーダ", "QA",
                   "テスト計画 / 実行統制",
                   "品質管理リーダ"],
                  ["独立検証 (IV&V)", "外部",
                   "重要機能・NFR の第三者検証",
                   "SteerCo"],
              ], font_size=11)


@builder("4", "品質管理計画", "4.3 品質管理プロセス")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "Plan → Do → Check → Act を各フェーズで反復し、品質ゲートで収束させる。",
             size=12, color=COLOR_GRAY)
    # PDCA 帯
    steps = [("Plan", "品質目標 / 指標設定"),
             ("Do",   "レビュー・テスト実施"),
             ("Check","メトリクス分析"),
             ("Act",  "是正・予防・再発防止")]
    x = Inches(0.5); y = Inches(2.4); w = Inches(3.05); h = Inches(1.1)
    colors = [COLOR_PRIMARY, COLOR_ACCENT, COLOR_PRIMARY, COLOR_ACCENT]
    for i, (h1, h2) in enumerate(steps):
        xi = x + (w + Inches(0.05)) * i
        add_rect(slide, xi, y, w, h, fill=colors[i])
        add_text(slide, xi, y + Inches(0.1), w, Inches(0.4), h1,
                 size=16, bold=True, color=COLOR_WHITE,
                 align=PP_ALIGN.CENTER)
        add_text(slide, xi, y + Inches(0.55), w, Inches(0.5), h2,
                 size=11, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    # ゲート基準
    add_text(slide, Inches(0.5), Inches(3.8), Inches(12), Inches(0.4),
             "品質ゲート基準 (例)", size=14, bold=True, color=COLOR_PRIMARY)
    add_table(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.3),
              headers=["工程", "レビュー指摘密度", "テスト密度", "欠陥密度", "残欠陥"],
              rows=[
                  ["設計",   "≥ 0.5 件/頁",  "-",              "-",          "Sev1/2 : 0"],
                  ["UT",     "-",           "≥ 1.0 ケース/KLOC", "≤ 0.5 件/KLOC", "Sev1/2 : 0"],
                  ["IT",     "-",           "計画比 100%",      "≤ 0.3 件/KLOC", "Sev1 : 0"],
                  ["ST/UAT", "-",           "計画比 100%",      "収束曲線上昇停止", "Sev1 : 0"],
              ], font_size=11)


@builder("4", "品質管理計画", "4.4 システム品質管理詳細")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "非機能要件 (NFR) を観点ごとに定量目標として管理する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["観点", "指標", "目標値", "検証方法"],
              rows=[
                  ["性能",     "応答時間 (95%tile)", "≤ 2.0 秒 (主要画面)", "負荷試験 (JMeter)"],
                  ["スループット", "トランザクション/秒", "≥ 200 TPS (ピーク時)", "負荷試験"],
                  ["可用性",   "月間稼働率",       "≥ 99.9% (8時-22時)", "監視 / 月次レポート"],
                  ["信頼性",   "Sev1 障害",       "年間 ≤ 2 件",         "障害記録"],
                  ["セキュリティ", "脆弱性",      "Critical/High : 0",   "SAST / DAST / 侵入試験"],
                  ["保守性",   "Cyclomatic 複雑度", "≤ 15 / 関数",        "SonarQube"],
                  ["使用性",   "UAT 満足度",      "≥ 4.0 / 5.0",         "アンケート"],
                  ["移行性",   "データ不整合率",  "≤ 0.01%",              "移行リハ結果"],
              ], font_size=11)


# ======================================================================
# 5. 移行計画
# ======================================================================
@builder("5", "移行計画", "5.1 リリース定義")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "段階リリースを採用。業務影響の小さい領域から順次展開する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["Release", "時期", "対象範囲", "方式", "並行運用"],
              rows=[
                  ["R0 (Pilot)", "2027-02",
                   "社内部門 (経理) のみ",
                   "限定ユーザ", "2 週間"],
                  ["R1 (Major)", "2027-03",
                   "販売 / 購買 / 在庫 / 会計",
                   "一斉切替 (Big Bang)", "1 ヶ月"],
                  ["R2",         "2027-06",
                   "人事 / 給与",
                   "月次切替", "-"],
                  ["R3 (BI)",    "2027-09",
                   "BI / ダッシュボード拡張",
                   "段階展開", "-"],
                  ["R4 (海外)",  "2027 下期〜 (Phase2)",
                   "海外拠点",
                   "別プロジェクト", "-"],
              ], font_size=11)


@builder("5", "移行計画", "5.2 リリース判定方針")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "定量・定性の両軸で Go / No-Go を判定する。1 項目でも NG なら原則 No-Go。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.0),
              headers=["判定項目", "基準", "判定者"],
              rows=[
                  ["品質",       "Sev1/2 欠陥 0、ST/UAT 完了率 100%", "QA"],
                  ["性能",       "NFR の性能目標を全て達成",         "アーキ + QA"],
                  ["移行",       "リハ 2 回連続成功、RTO/RPO 達成",    "移行リーダ"],
                  ["運用",       "手順書整備、運用リハ完了、要員訓練済", "運用責任者"],
                  ["業務",       "UAT 承認、教育完了、業務部門 Go",     "業務オーナ"],
                  ["BCP",        "切戻し手順の実機確認完了",           "PM + 運用"],
              ], font_size=11)
    add_text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
             "最終判定会議: リリース予定日の 7 営業日前に SteerCo で実施。",
             size=11, bold=True, color=COLOR_ACCENT)


@builder("5", "移行計画", "5.3 システム保守計画")
def _(slide):
    col_w = Inches(6.1); y = Inches(1.9); h = Inches(5.0)
    add_box(slide, Inches(0.5), y, col_w, h, "保守体制 (ハイパーケア → 定常)", [
        ("ハイパーケア", "本稼働後 6 ヶ月。PJ メンバ常駐・即応"),
        ("定常運用",     "7 ヶ月目以降。運用保守チームへ移管"),
        ("体制",         "1次: ヘルプデスク / 2次: 運用 / 3次: 開発"),
        ("稼働時間",     "平日 8:00-22:00、夜間・休日はオンコール"),
        ("SLA",          "受付 < 15分 / 初動 < 1h / 復旧 < 4h (Sev1)"),
    ])
    add_box(slide, Inches(6.7), y, col_w, h, "保守区分と対応範囲",
            [
                ("是正保守", "障害対応 (SLA に基づく復旧)"),
                ("予防保守", "パッチ適用 / 脆弱性対応 / 監視改善"),
                ("適応保守", "OS・ミドルアップデート対応"),
                ("完全化保守", "小規模改善 (四半期ごとリリース)"),
                ("機能追加", "別契約 / 別プロジェクトで対応"),
                ("契約",    "年次更新 / SLA 未達はクレジット発行"),
            ], title_color=COLOR_ACCENT)


@builder("5", "移行計画", "5.4 業務運用計画")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "業務部門が自律運用に移れる状態を移行完了の定義とする。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["領域", "内容", "成果物", "責任部門"],
              rows=[
                  ["業務手順", "新業務フロー整備・手順書化",
                   "業務マニュアル (版管理)", "業務部門"],
                  ["教育",     "役割別トレーニング / e-Learning",
                   "教育計画 / 受講記録",     "人事 + 業務"],
                  ["ヘルプデスク", "問合せ窓口・FAQ 整備",
                   "FAQ / ナレッジ DB",       "情シス"],
                  ["KPI モニタ", "KGI/KPI の測定と経営報告",
                   "月次ダッシュボード",      "経営企画"],
                  ["継続改善",   "現場フィードバックの収集・改善",
                   "改善要望一覧",            "PMO → 運用"],
                  ["BCP",        "災害 / 障害時の業務継続手順",
                   "BCP マニュアル",          "リスク管理"],
              ], font_size=11)


# ======================================================================
# 6. AI活用計画 (Kiro)
# ======================================================================
@builder("6", "AI活用計画", "6.1 AI活用方針")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "生成AIを単なるコーディング補助ではなく、仕様駆動開発の中核に据える。",
             size=12, color=COLOR_GRAY)
    items = [
        ("方針1", "Spec-Driven Development — 要件・設計・タスクを AI と共同で形式化し、実装の前段を厚くする"),
        ("方針2", "Human-in-the-Loop — 仕様・設計・マージは必ず人が承認。AI の自律実行は定義済み範囲のみ"),
        ("方針3", "再現性 — プロンプト / spec / hooks / steering をリポジトリで版管理し、属人性を排除"),
        ("方針4", "セキュリティ優先 — 機密データは社外 API に出さない。MCP 経由で社内リソース接続"),
        ("方針5", "効果測定 — 生産性・品質・リードタイムを定量計測し、PDCA を回す"),
    ]
    add_bullets(slide, Inches(0.8), Inches(2.4), Inches(11.8), Inches(3.3),
                items, size=13, line_spacing=1.45)
    # 適用領域
    add_text(slide, Inches(0.5), Inches(5.8), Inches(12), Inches(0.4),
             "適用領域", size=14, bold=True, color=COLOR_PRIMARY)
    cats = ["要件定義", "設計", "実装", "テスト", "レビュー", "ドキュメント", "運用/保守"]
    x = Inches(0.5); y = Inches(6.25); w = Inches(1.77); h = Inches(0.6)
    for i, c in enumerate(cats):
        add_rect(slide, x + (w + Inches(0.05)) * i, y, w, h,
                 fill=COLOR_LIGHT, line=COLOR_ACCENT)
        add_text(slide, x + (w + Inches(0.05)) * i, y, w, h, c,
                 size=11, bold=True, color=COLOR_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


@builder("6", "AI活用計画", "6.2 Kiro の活用 (Specs / Hooks / Steering / MCP)")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "Kiro (Agentic IDE) の 4 つのコア機能を本プロジェクトに組み込む。",
             size=12, color=COLOR_GRAY)
    # 2x2 grid
    x0 = Inches(0.5); y0 = Inches(2.3)
    w = Inches(6.15); h = Inches(2.25); gx = Inches(0.2); gy = Inches(0.2)
    # Specs
    add_box(slide, x0, y0, w, h,
            "① Specs (仕様駆動)",
            [
                ("requirements.md", "EARS 記法で受入条件を列挙"),
                ("design.md",       "アーキ / IF / データモデル / 意思決定"),
                ("tasks.md",        "実装タスクへの分解 (トレーサビリティ確保)"),
                ("成果",            "要件 → 実装の追跡可能性を担保"),
            ], title_color=COLOR_PRIMARY)
    # Hooks
    add_box(slide, x0 + w + gx, y0, w, h,
            "② Agent Hooks (自動化)",
            [
                ("on save",         "該当テストの自動実行 / Lint / 型検査"),
                ("on commit",       "コミットメッセージ整形 / 変更要約生成"),
                ("on spec update",  "影響範囲のタスク再生成"),
                ("成果",            "反復作業を AI が肩代わり"),
            ], title_color=COLOR_ACCENT)
    # Steering
    add_box(slide, x0, y0 + h + gy, w, h,
            "③ Steering (文脈制御)",
            [
                ("product.md",      "プロダクト方針 / ペルソナ / KPI"),
                ("tech.md",         "技術スタック / コーディング規約"),
                ("structure.md",    "リポジトリ構造 / 命名規則"),
                ("成果",            "AI 出力の一貫性と品質を底上げ"),
            ], title_color=COLOR_PRIMARY)
    # MCP
    add_box(slide, x0 + w + gx, y0 + h + gy, w, h,
            "④ MCP (社内連携)",
            [
                ("社内DB / Wiki",   "読み取り専用で Kiro から参照"),
                ("Jira / GitHub",   "チケット起票・PR 操作を AI から"),
                ("監視 / ログ",     "障害分析時のエビデンス取得"),
                ("成果",            "機密を外に出さず社内情報で推論"),
            ], title_color=COLOR_ACCENT)


@builder("6", "AI活用計画", "6.3 開発プロセスへの統合")
def _(slide):
    add_text(slide, Inches(0.5), Inches(1.8), Inches(12), Inches(0.4),
             "各フェーズでの AI 活用ポイントと、人が担う責任領域を明示する。",
             size=12, color=COLOR_GRAY)
    add_table(slide, Inches(0.5), Inches(2.3), Inches(12.3), Inches(4.5),
              headers=["フェーズ", "AI が担う作業 (Kiro)", "人が担う責任",
                       "指標"],
              rows=[
                  ["要件定義",
                   "spec 雛形生成 / 矛盾検出 / 用語統一",
                   "業務要件の妥当性判断 / 承認",
                   "要件レビュー時間 -30%"],
                  ["基本設計",
                   "代替案提示 / IF 差分検出 / NFR チェック",
                   "アーキ決定 / ADR 承認",
                   "設計指摘 再発率 -20%"],
                  ["詳細設計/実装",
                   "コード生成 / リファクタ / テスト雛形",
                   "設計意図の判断 / レビュー / マージ",
                   "開発生産性 +30%"],
                  ["テスト",
                   "テストケース補完 / データ生成 / 実行",
                   "受入観点・業務妥当性の確認",
                   "UT カバレッジ ≥ 80%"],
                  ["レビュー",
                   "静的解析 / 規約違反検出 / 差分要約",
                   "設計意図・ビジネス妥当性の判断",
                   "レビュー待ち時間 -40%"],
                  ["ドキュメント",
                   "API 仕様 / 運用手順の自動生成と更新",
                   "最終校正・顧客向け表現の調整",
                   "ドキュ更新遅延 0 日"],
                  ["運用/保守",
                   "障害一次切り分け / 再発防止案提示",
                   "復旧判断 / 顧客コミュニケーション",
                   "MTTR -25%"],
              ], font_size=10)


@builder("6", "AI活用計画", "6.4 AIガバナンス (セキュリティ・品質・責任)")
def _(slide):
    col_w = Inches(6.1); y = Inches(1.9); h = Inches(5.0)
    add_box(slide, Inches(0.5), y, col_w, h,
            "セキュリティ / プライバシー",
            [
                ("データ境界", "顧客データ・個人情報は AI に投入禁止 (DLP 監視)"),
                ("モデル",    "社内承認モデル限定 / ログは国内リージョン保管"),
                ("秘密情報", "API Key / 認証情報のプロンプト投入を自動検出"),
                ("MCP",      "最小権限・読み取り優先 / 書込は承認フロー経由"),
                ("ログ監査", "プロンプト / 生成物 / 参照コンテキストを1年保管"),
            ])
    add_box(slide, Inches(6.7), y, col_w, h,
            "品質 / 責任 / 運用",
            [
                ("責任原則", "AI 生成物の最終責任は常に人 (レビュア)"),
                ("検証",    "生成コードは CI (Lint/型/テスト/SAST) 必須通過"),
                ("著作権",  "ライセンス互換性チェック / 類似コード検出"),
                ("幻覚対策", "根拠 (Spec/コード/ドキュ) 提示を徹底"),
                ("効果測定", "月次で生産性・欠陥密度・レビュー時間を報告"),
                ("教育",    "プロンプト設計・Spec 記法の研修を全員必修"),
            ], title_color=COLOR_ACCENT)


# ======================================================================
# 最終ページ
# ======================================================================
def build_end(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=COLOR_PRIMARY)
    add_text(slide, 0, Inches(3.0), SLIDE_W, Inches(1.2),
             "Thank You", size=60, bold=True, color=COLOR_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(slide, 0, Inches(4.2), SLIDE_W, Inches(0.6),
             "ご質問・ご意見はプロジェクト窓口まで", size=16,
             color=COLOR_LIGHT, align=PP_ALIGN.CENTER)


# ---- ビルド ----------------------------------------------------------
# 全スライド数 = 表紙 + 目次 + (章扉 + 各サブページ)×5章 + Thanks
section_map = {}
for sec_no, sec_title, _, _ in PAGES:
    section_map.setdefault(sec_no, sec_title)

total_slides = 2  # cover + toc
for sec_no in section_map:
    total_slides += 1  # divider
    total_slides += sum(1 for p in PAGES if p[0] == sec_no)
total_slides += 1  # thanks

# 表紙
s = prs.slides.add_slide(BLANK)
build_cover(s)

# 目次
s = prs.slides.add_slide(BLANK)
build_toc(s)

# ページ順に: 章扉 → 各サブページ
page_no = 2  # 表紙=1, 目次=2
section_subtitles = {
    "1": "Project Definition",
    "2": "Project Organization",
    "3": "Project Management Plan",
    "4": "Quality Management Plan",
    "5": "Migration Plan",
    "6": "AI Utilization Plan — powered by Kiro",
}
current_section = None
for sec_no, sec_title, page_title, fn in PAGES:
    if sec_no != current_section:
        # 章扉
        s = prs.slides.add_slide(BLANK)
        section_divider(sec_no, sec_title,
                        section_subtitles.get(sec_no, ""))(s)
        page_no += 1
        current_section = sec_no
    s = prs.slides.add_slide(BLANK)
    add_page_header(s, sec_no, sec_title, page_title)
    fn(s)
    page_no += 1
    add_footer(s, page_no, total_slides)

# Thanks
s = prs.slides.add_slide(BLANK)
build_end(s)

# ============================================================================
#   出力処理
# ============================================================================
# OUTPUT_FILE が相対パスならスクリプトと同じフォルダに出力。
# 絶対パスならそのまま使用。
def _resolve_output_path(filename: str) -> Path:
    """出力パスを解決する。

    - 絶対パス指定: そのまま使用
    - 相対パス指定: このスクリプトと同じフォルダ直下に作成
      (対話実行 / Jupyter などで __file__ が無い場合は cwd にフォールバック)
    """
    p = Path(filename).expanduser()
    if p.is_absolute():
        return p
    try:
        base = Path(__file__).resolve().parent
    except NameError:
        base = Path.cwd()
    return base / p


output_path = _resolve_output_path(OUTPUT_FILE)
output_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(output_path))

print("=" * 60)
print(f"  OK: プロジェクト計画書を生成しました")
print(f"  path : {output_path}")
print(f"  slides: {len(prs.slides)} pages")
print("=" * 60)
print("  PowerPoint / Keynote / LibreOffice Impress で開けます。")
